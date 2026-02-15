"""
Generate IntelliWheels Final Pitch as PPTX.
Run: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
PURPLE = RGBColor(0x7B, 0x2F, 0xF7)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
GREEN = RGBColor(0x25, 0xD3, 0x66)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
RED = RGBColor(0xFF, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)


def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_text(text_frame, text, font_size=13, color=LIGHT_GRAY, bold_prefix=None):
    p = text_frame.add_paragraph()
    if bold_prefix:
        run = p.add_run()
        run.text = f"{bold_prefix}: "
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE
        run.font.name = 'Calibri'
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = 'Calibri'
    else:
        p.text = f"▸ {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
    p.space_before = Pt(3)
    p.space_after = Pt(3)
    return p


def card_with_text(slide, left, top, width, height, icon, title, body, border_color=None):
    shape = add_shape(slide, left, top, width, height, BG_CARD, border_color or PURPLE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.space_before = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(6)


def add_tag(slide, left, top, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.6), Inches(0.35))
    shape.fill.solid()
    # Create a darker version of the color for background
    hex_str = str(color)
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    shape.fill.fore_color.rgb = RGBColor(r // 4, g // 4, b // 4)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ================================================================
# SLIDE 1: TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, Inches(0), Inches(0.8), Inches(13.333), Inches(0.8),
             "\U0001F697", 60, WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.8), Inches(13.333), Inches(1),
             "IntelliWheels", 48, PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.9), Inches(13.333), Inches(0.6),
             "AI-Powered Car Marketplace Assistant for Jordan", 20, CYAN, alignment=PP_ALIGN.CENTER)

tag_y = Inches(3.8)
tag_start = Inches(2.8)
tags = [("Gemini 2.5", PURPLE), ("Django REST", CYAN), ("Twilio WhatsApp", GREEN),
        ("Crisp Logic", ORANGE), ("Vision AI", RED)]
for i, (txt, clr) in enumerate(tags):
    add_tag(slide, tag_start + Inches(i * 1.7), tag_y, txt, clr)

add_text_box(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.4),
             "AI Engineer Course — Capstone Project — February 2026", 12, DARK_GRAY,
             alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 2: PROBLEM
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "The Problem", 32, ORANGE, bold=True)
add_text_box(slide, Inches(0.7), Inches(1.0), Inches(8), Inches(0.4),
             "Car buyers in Jordan face multiple challenges", 14, GRAY)

# 4 problem cards
cards = [
    ("💰", "Price Opacity", "No standardized reference for used car\nvaluations in JOD. Sellers set arbitrary\nprices — buyers can't verify.", RED),
    ("🔍", "Information Fragmentation", "Specs, comparisons, and advice scattered\nacross forums, Facebook groups,\nand word of mouth.", ORANGE),
    ("📷", "No Visual Identification", "Buyers see a car on the street and can't\nquickly identify what it is, what year,\nor what it's worth.", PURPLE),
    ("📱", "Accessibility Gap", "Technical car data isn't easily accessible\nto non-technical users, especially\non mobile devices.", CYAN),
]
for i, (icon, title, body, clr) in enumerate(cards):
    col = i % 2
    row = i // 2
    card_with_text(slide, Inches(0.7 + col * 6.1), Inches(1.7 + row * 2.7),
                   Inches(5.8), Inches(2.3), icon, title, body, clr)

# ================================================================
# SLIDE 3: SOLUTION
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "The Solution: IntelliWheels", 32, CYAN, bold=True)
add_text_box(slide, Inches(0.7), Inches(1.0), Inches(8), Inches(0.4),
             "4 AI-powered features in one platform", 14, GRAY)

solutions = [
    ("🤖", "AI Chatbot", "Google Gemini 2.5 Flash — Ask anything\nabout cars: specs, comparisons,\nrecommendations. Jordanian market\nfocus with JOD pricing.", PURPLE),
    ("📊", "Price Estimator", "Crisp (rule-based) logic — 15 makes,\n70+ models. Three depreciation tiers:\nLuxury (12%), Premium (10%),\nEconomic (7%).", ORANGE),
    ("👁️", "Vision Helper", "Gemini Vision API — Upload any car\nphoto and instantly identify make,\nmodel, year, and visible condition.", CYAN),
    ("📲", "WhatsApp Share", "Twilio API — Send chat replies and\nprice estimates directly to WhatsApp\nfor on-the-go access.", GREEN),
]
for i, (icon, title, body, clr) in enumerate(solutions):
    col = i % 2
    row = i // 2
    card_with_text(slide, Inches(0.7 + col * 6.1), Inches(1.7 + row * 2.7),
                   Inches(5.8), Inches(2.3), icon, title, body, clr)

# ================================================================
# SLIDE 4: TECH STACK
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "Technology Stack", 32, CYAN, bold=True)

# Tech table as text
tech_shape = add_shape(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, PURPLE)
tf = tech_shape.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""

tech_rows = [
    ("Language", "Python 3.12"),
    ("Backend", "Django 4.2 + REST Framework"),
    ("LLM", "Google Gemini 2.5 Flash"),
    ("LLM SDK", "google-genai (new official SDK)"),
    ("Vision", "Gemini Vision API"),
    ("Price Logic", "Crisp / rule-based (custom)"),
    ("Messaging", "Twilio WhatsApp Sandbox"),
    ("Database", "SQLite 3"),
    ("Frontend", "HTML / CSS / JavaScript"),
    ("Config", "python-dotenv (.env)"),
]
for label, val in tech_rows:
    p = tf.add_paragraph()
    run1 = p.add_run()
    run1.text = f"{label}:  "
    run1.font.size = Pt(13)
    run1.font.color.rgb = CYAN
    run1.font.bold = True
    run1.font.name = 'Calibri'
    run2 = p.add_run()
    run2.text = val
    run2.font.size = Pt(13)
    run2.font.color.rgb = WHITE
    run2.font.name = 'Calibri'
    p.space_before = Pt(6)
    p.space_after = Pt(2)

# Why these choices
why_box = add_text_box(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5),
                       "Why These Choices?", 18, WHITE, bold=True)
reasons_box = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(4.8))
tf2 = reasons_box.text_frame
tf2.word_wrap = True
tf2.paragraphs[0].text = ""
add_bullet_text(tf2, "Battle-tested framework with admin, ORM, serializer validation", bold_prefix="Django + DRF")
add_bullet_text(tf2, "Fast, multimodal, free tier suitable for development", bold_prefix="Gemini 2.5 Flash")
add_bullet_text(tf2, "New official SDK replacing deprecated google-generativeai", bold_prefix="google-genai SDK")
add_bullet_text(tf2, "Deterministic pricing — no hallucinated numbers", bold_prefix="Crisp logic")
add_bullet_text(tf2, "Industry-standard WhatsApp API with sandbox for prototyping", bold_prefix="Twilio")

# ================================================================
# SLIDE 5: ARCHITECTURE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "System Architecture", 32, CYAN, bold=True)

# Architecture boxes
layers = [
    (Inches(2), Inches(1.3), Inches(9), Inches(0.7), "🌐 Browser — Dark Cyberpunk UI (Single Page)", CYAN),
    (Inches(2), Inches(2.3), Inches(9), Inches(0.7), "🚀 Django REST Framework — 7 API Endpoints", PURPLE),
]
for l, t, w, h, txt, clr in layers:
    shape = add_shape(slide, l, t, w, h, BG_CARD, clr)
    shape.text_frame.paragraphs[0].text = txt
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = 'Calibri'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Arrow
add_text_box(slide, Inches(6), Inches(2.0), Inches(1), Inches(0.4), "⬇", 16, CYAN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6), Inches(3.0), Inches(1), Inches(0.4), "⬇", 16, CYAN, alignment=PP_ALIGN.CENTER)

# 4 engine boxes
engines = [
    ("🤖 Car\nChatbot", PURPLE), ("📊 Crisp\nLogic", ORANGE),
    ("👁️ Gemini\nVision", CYAN), ("📲 Twilio\nWhatsApp", GREEN),
]
for i, (txt, clr) in enumerate(engines):
    shape = add_shape(slide, Inches(2 + i * 2.35), Inches(3.3), Inches(2.1), Inches(1.1), BG_CARD, clr)
    shape.text_frame.paragraphs[0].text = txt
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = 'Calibri'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text_box(slide, Inches(6), Inches(4.4), Inches(1), Inches(0.4), "⬇", 16, CYAN, alignment=PP_ALIGN.CENTER)

# Bottom boxes
for i, (txt, clr) in enumerate([("🗄️ SQLite — All records persisted", GRAY), ("🧠 Gemini API (google-genai)", PURPLE)]):
    shape = add_shape(slide, Inches(2 + i * 4.7), Inches(4.8), Inches(4.4), Inches(0.7), BG_CARD, clr)
    shape.text_frame.paragraphs[0].text = txt
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = 'Calibri'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Data flows on the right
flows_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1.5))
tf = flows_box.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""
add_bullet_text(tf, "User → ChatView → CarChatbotAgent → GeminiClient → Gemini API → ChatMessage → JSON", 11, GRAY, "Chat")
add_bullet_text(tf, "User → PriceEstimateView → estimate_price() crisp logic → PriceEstimate → JSON", 11, GRAY, "Price")
add_bullet_text(tf, "Image → VisionView → analyze_car_image() → Gemini Vision → VisionAnalysis → JSON", 11, GRAY, "Vision")
add_bullet_text(tf, "Click share → WhatsAppView → send_whatsapp_message() → Twilio API → Delivered", 11, GRAY, "WhatsApp")

# ================================================================
# SLIDE 6: API ENDPOINTS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "API Design — 7 Endpoints", 32, CYAN, bold=True)

endpoints = [
    ("/api/chat/", "POST", "AI chatbot conversation", "message, session_id (opt)"),
    ("/api/estimate/", "POST", "Crisp logic price estimate", "make, model, year, mileage"),
    ("/api/makes/", "GET", "List all 15 car makes", "—"),
    ("/api/models/<make>/", "GET", "List models for a make", "make (URL param)"),
    ("/api/vision/", "POST", "Car image AI analysis", "image (file upload)"),
    ("/api/whatsapp/send/", "POST", "Send via WhatsApp", "message, to (opt)"),
    ("/health/", "GET", "Health check", "—"),
]

# Table header
header_shape = add_shape(slide, Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.5), RGBColor(0x24, 0x1B, 0x4A))
tf = header_shape.text_frame
tf.paragraphs[0].text = ""
p = tf.paragraphs[0]
cols = ["Endpoint", "Method", "Description", "Input"]
p.text = f"{'Endpoint':<28} {'Method':<10} {'Description':<30} {'Input'}"
p.font.size = Pt(12)
p.font.color.rgb = CYAN
p.font.bold = True
p.font.name = 'Consolas'

for i, (ep, method, desc, inp) in enumerate(endpoints):
    row_shape = add_shape(slide, Inches(0.7), Inches(1.75 + i * 0.52), Inches(11.9), Inches(0.48),
                          BG_CARD if i % 2 == 0 else RGBColor(0x14, 0x14, 0x28))
    row_shape.line.fill.background()
    tf = row_shape.text_frame
    tf.paragraphs[0].text = f"{ep:<28} {method:<10} {desc:<30} {inp}"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = 'Consolas'
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Example request/response
add_text_box(slide, Inches(0.7), Inches(5.6), Inches(5.8), Inches(0.3),
             "Example: POST /api/estimate/", 12, CYAN, bold=True)
code1 = add_shape(slide, Inches(0.7), Inches(5.95), Inches(5.8), Inches(1.2), RGBColor(0x0D, 0x0D, 0x18), GRAY)
code1.text_frame.paragraphs[0].text = '{ "make": "toyota", "model": "camry",\n  "year": 2020, "mileage_km": 80000 }\n\n→ { "depreciated_price_jod": 16213, "depreciation_pct": 35.1 }'
code1.text_frame.paragraphs[0].font.size = Pt(10)
code1.text_frame.paragraphs[0].font.color.rgb = CYAN
code1.text_frame.paragraphs[0].font.name = 'Consolas'

add_text_box(slide, Inches(7), Inches(5.6), Inches(5.8), Inches(0.3),
             "Example: POST /api/chat/", 12, CYAN, bold=True)
code2 = add_shape(slide, Inches(7), Inches(5.95), Inches(5.8), Inches(1.2), RGBColor(0x0D, 0x0D, 0x18), GRAY)
code2.text_frame.paragraphs[0].text = '{ "message": "Best car under 15K JOD?" }\n\n→ { "message": "For 15,000 JOD...",\n    "model_used": "gemini-2.5-flash" }'
code2.text_frame.paragraphs[0].font.size = Pt(10)
code2.text_frame.paragraphs[0].font.color.rgb = CYAN
code2.text_frame.paragraphs[0].font.name = 'Consolas'

# ================================================================
# SLIDE 7: PRICE ESTIMATOR
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "Crisp Price Estimator", 32, ORANGE, bold=True)
add_text_box(slide, Inches(0.7), Inches(1.0), Inches(10), Inches(0.4),
             "Deterministic, rule-based — no hallucinated prices", 14, GRAY)

# Formula
formula = add_shape(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(0.8), RGBColor(0x0D, 0x0D, 0x18), ORANGE)
formula.text_frame.paragraphs[0].text = "Value = Base Price × (1 − rate)^age − Mileage Penalty"
formula.text_frame.paragraphs[0].font.size = Pt(16)
formula.text_frame.paragraphs[0].font.color.rgb = WHITE
formula.text_frame.paragraphs[0].font.bold = True
formula.text_frame.paragraphs[0].font.name = 'Calibri'
formula.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
formula.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Categories table
cats = [
    ("LUXURY", "12%/yr", "S-Class, 7-Series, Cayenne", RED),
    ("PREMIUM", "10%/yr", "C-Class, 3-Series, A4", ORANGE),
    ("ECONOMIC", "7%/yr", "Camry, Elantra, Civic", CYAN),
]
for i, (cat, rate, ex, clr) in enumerate(cats):
    shape = add_shape(slide, Inches(0.7), Inches(2.6 + i * 0.7), Inches(5.8), Inches(0.6), BG_CARD, clr)
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"  {cat}  "
    run1.font.size = Pt(12)
    run1.font.color.rgb = clr
    run1.font.bold = True
    run1.font.name = 'Calibri'
    run2 = p.add_run()
    run2.text = f"  {rate}    {ex}"
    run2.font.size = Pt(12)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = 'Calibri'

# Mileage penalties
mile_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.9), Inches(5.8), Inches(1.5))
tf = mile_box.text_frame
tf.word_wrap = True
add_paragraph(tf, "Mileage Penalties", 16, WHITE, bold=True)
add_bullet_text(tf, "> 100,000 km → extra 5% deduction")
add_bullet_text(tf, "> 200,000 km → extra 10% deduction")
add_bullet_text(tf, "Floor: never below 5% of base price")

# Stats
stats = [("15", "Car Makes"), ("70+", "Models"), ("3", "Categories")]
stat_colors = [CYAN, ORANGE, GREEN]
for i, (num, label) in enumerate(stats):
    add_text_box(slide, Inches(7.2 + i * 2), Inches(1.6), Inches(1.8), Inches(0.6),
                 num, 36, stat_colors[i], bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.2 + i * 2), Inches(2.2), Inches(1.8), Inches(0.3),
                 label, 12, GRAY, alignment=PP_ALIGN.CENTER)

# Example calc
ex_shape = add_shape(slide, Inches(7), Inches(2.8), Inches(5.8), Inches(3.8), BG_CARD, GREEN)
tf = ex_shape.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Example Calculation"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'

add_paragraph(tf, "Toyota Camry 2020 • 80,000 km", 14, WHITE, bold=True, space_before=Pt(12))
lines = [
    "Category: Economic (7%/yr)",
    "Base: 25,000 JOD",
    "Age: 6 years → factor: 0.6485",
    "After age: 16,213 JOD",
    "Mileage: under 100K → no penalty",
]
for line in lines:
    add_paragraph(tf, line, 12, LIGHT_GRAY)
add_paragraph(tf, "Final: 16,213 JOD (35.1% depreciation)", 14, GREEN, bold=True, space_before=Pt(8))

# ================================================================
# SLIDE 8: LLM & AGENT
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "LLM Integration & Agent", 32, CYAN, bold=True)

# GeminiClient card
gc_shape = add_shape(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(2.5), BG_CARD, PURPLE)
tf = gc_shape.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "GeminiClient Architecture"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'

details = [
    ("SDK", "google-genai (new official SDK)"),
    ("Model", "gemini-2.5-flash"),
    ("Pattern", "Singleton instance (one client, all requests)"),
    ("Methods", "chat() + analyze_image()"),
]
for label, val in details:
    p = tf.add_paragraph()
    r1 = p.add_run()
    r1.text = f"{label}: "
    r1.font.size = Pt(12)
    r1.font.color.rgb = PURPLE
    r1.font.bold = True
    r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(12)
    r2.font.color.rgb = WHITE
    r2.font.name = 'Calibri'
    p.space_before = Pt(4)

# Key features
feat_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(5.8), Inches(3.2))
tf2 = feat_box.text_frame
tf2.word_wrap = True
add_paragraph(tf2, "Key Features", 16, WHITE, bold=True)
add_bullet_text(tf2, "Jordanian market specialization via GenerateContentConfig", bold_prefix="System instructions")
add_bullet_text(tf2, "Full conversation history with Content/Part objects", bold_prefix="Multi-turn")
add_bullet_text(tf2, "Part.from_bytes() for image analysis", bold_prefix="Vision")
add_bullet_text(tf2, "Works without API key (returns setup instructions)", bold_prefix="Mock mode")
add_bullet_text(tf2, "429 errors → friendly message, no crashes", bold_prefix="Rate limit handling")

# Agent card on right
agent_shape = add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(3.2), BG_CARD, PURPLE)
tf3 = agent_shape.text_frame
tf3.word_wrap = True
tf3.paragraphs[0].text = "CarChatbotAgent — System Prompt"
tf3.paragraphs[0].font.size = Pt(16)
tf3.paragraphs[0].font.color.rgb = WHITE
tf3.paragraphs[0].font.bold = True
tf3.paragraphs[0].font.name = 'Calibri'

capabilities = [
    "Car specs & comparisons",
    "Budget recommendations (JOD)",
    "Category explanations (Luxury/Premium/Economic)",
    "Price guidance → redirect to estimator",
    "Maintenance & import tips for Jordan",
    "Politely redirects non-car questions",
]
for cap in capabilities:
    add_paragraph(tf3, f"▸ {cap}", 12, LIGHT_GRAY, space_before=Pt(4))

# Session mgmt
sess_shape = add_shape(slide, Inches(7), Inches(4.8), Inches(5.8), Inches(2.2), BG_CARD, CYAN)
tf4 = sess_shape.text_frame
tf4.word_wrap = True
tf4.paragraphs[0].text = "Session Management"
tf4.paragraphs[0].font.size = Pt(16)
tf4.paragraphs[0].font.color.rgb = WHITE
tf4.paragraphs[0].font.bold = True
tf4.paragraphs[0].font.name = 'Calibri'

sess_items = [
    "Each conversation → unique ChatSession",
    "Every message saved as ChatMessage",
    "Full history sent to Gemini for context",
    "Frontend tracks session_id for continuity",
]
for item in sess_items:
    add_paragraph(tf4, f"▸ {item}", 12, LIGHT_GRAY, space_before=Pt(3))

# ================================================================
# SLIDE 9: VISION & WHATSAPP
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "Vision Analysis & WhatsApp", 32, GREEN, bold=True)

# Vision card
vis_shape = add_shape(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, CYAN)
tf = vis_shape.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "👁️ Gemini Vision Helper"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'

vis_details = [("Input", "Car photo (JPEG/PNG/WebP)"), ("Output", "Make, Model, Year, Condition"),
               ("Method", "Part.from_bytes() → Gemini Vision"), ("Prompt", "Structured JSON output format")]
for label, val in vis_details:
    p = tf.add_paragraph()
    r1 = p.add_run()
    r1.text = f"{label}: "
    r1.font.size = Pt(12)
    r1.font.color.rgb = CYAN
    r1.font.bold = True
    r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(12)
    r2.font.color.rgb = WHITE
    r2.font.name = 'Calibri'
    p.space_before = Pt(4)

# Vision parsing
vis2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(5.8), Inches(2.5))
tf2 = vis2.text_frame
tf2.word_wrap = True
add_paragraph(tf2, "Response Parsing", 14, WHITE, bold=True)
add_bullet_text(tf2, "Parse JSON from Gemini response")
add_bullet_text(tf2, "Strip markdown code fences if present")
add_bullet_text(tf2, "Fallback: raw text → condition field")
add_bullet_text(tf2, "All results saved to VisionAnalysis model")

# WhatsApp card
wa_shape = add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(2.8), BG_CARD, GREEN)
tf3 = wa_shape.text_frame
tf3.word_wrap = True
tf3.paragraphs[0].text = "📲 Twilio WhatsApp Integration"
tf3.paragraphs[0].font.size = Pt(18)
tf3.paragraphs[0].font.color.rgb = WHITE
tf3.paragraphs[0].font.bold = True
tf3.paragraphs[0].font.name = 'Calibri'

wa_details = [("Provider", "Twilio WhatsApp Sandbox"), ("Pattern", "Follows agents_htu-main project"),
              ("Sender", "+14155238886 (sandbox)"), ("Recipient", "Configurable via .env")]
for label, val in wa_details:
    p = tf3.add_paragraph()
    r1 = p.add_run()
    r1.text = f"{label}: "
    r1.font.size = Pt(12)
    r1.font.color.rgb = GREEN
    r1.font.bold = True
    r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(12)
    r2.font.color.rgb = WHITE
    r2.font.name = 'Calibri'
    p.space_before = Pt(4)

# WhatsApp features
wa2 = slide.shapes.add_textbox(Inches(7), Inches(4.3), Inches(5.8), Inches(2.5))
tf4 = wa2.text_frame
tf4.word_wrap = True
add_paragraph(tf4, "Features", 14, WHITE, bold=True)
add_bullet_text(tf4, "Auto-add whatsapp: prefix to numbers")
add_bullet_text(tf4, "Share chat responses with one click")
add_bullet_text(tf4, "Share price estimates with one click")
add_bullet_text(tf4, "Error 63007 handling (invalid sender)")
add_bullet_text(tf4, "Toast notifications for send status")

# ================================================================
# SLIDE 10: DATABASE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "Database Schema — 5 Models", 32, CYAN, bold=True)

models_data = [
    ("🚗", "CarListing", "make, model, year, mileage_km,\nfuel_type, transmission, category,\ncolor, prices, description, is_sold", PURPLE),
    ("💬", "ChatSession", "started_at, ended_at, summary\n\n← has many ChatMessages", CYAN),
    ("💬", "ChatMessage", "session (FK), role (user/assistant),\ncontent, created_at\n\nFull conversation history", CYAN),
    ("📊", "PriceEstimate", "make, model, year, mileage_km,\ncategory, original_price_jod,\ndepreciated_price, breakdown (JSON)", ORANGE),
    ("👁️", "VisionAnalysis", "image (file), detected_make,\ndetected_model, detected_year,\ncondition_summary, raw_response", GREEN),
    ("⚙️", "Admin Panel", "All 5 models registered with\nlist_display, filters, and\nsearch fields at /admin/", GRAY),
]
for i, (icon, title, body, clr) in enumerate(models_data):
    col = i % 3
    row = i // 3
    card_with_text(slide, Inches(0.7 + col * 4.1), Inches(1.3 + row * 3.0),
                   Inches(3.8), Inches(2.7), icon, title, body, clr)

add_text_box(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.4),
             "SQLite database  •  Django ORM  •  All records persisted and inspectable via /admin/",
             11, DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 11: SAFETY & PERFORMANCE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "Safety & Performance", 32, RED, bold=True)

safety_cards = [
    ("🔒", "Input Validation", "DRF serializers validate all inputs.\nImage MIME whitelist. Year 1980–2026.\nMessage max 5000 chars.", RED),
    ("⚡", "Rate Limit Handling", "Gemini 429/RESOURCE_EXHAUSTED\ncaught → friendly message.\nNo crashes on quota exhaustion.", ORANGE),
    ("📋", "Structured Logging", "3 loggers (core, api, ai_engine).\napp.log + errors.log.\nVerbose format with PID/thread.", CYAN),
    ("🛡️", "CORS & CSRF", "CORS restricted in production.\nCSRF middleware active.\nAuth validators enabled.", PURPLE),
    ("🔄", "Graceful Fallbacks", "Mock responses (no API key).\nVision JSON parse fallback.\nFloor price (5% of base).", GREEN),
    ("📦", "Data Persistence", "All chats, estimates, and vision\nresults saved. Django Admin\nfor full data inspection.", GRAY),
]
for i, (icon, title, body, clr) in enumerate(safety_cards):
    col = i % 3
    row = i // 3
    card_with_text(slide, Inches(0.7 + col * 4.1), Inches(1.3 + row * 3.0),
                   Inches(3.8), Inches(2.6), icon, title, body, clr)

# ================================================================
# SLIDE 12: CAPSTONE CHECKLIST
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
             "Capstone Requirements  ✓", 32, GREEN, bold=True)

# Core requirements
core_shape = add_shape(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, GREEN)
tf = core_shape.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Core Requirements"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Calibri'

core_items = [
    "✅ Use case design (car marketplace for Jordan)",
    "✅ Django project + API endpoints (7 endpoints)",
    "✅ Core LLM logic (Gemini via google-genai SDK)",
    "✅ Agent capabilities (CarChatbotAgent)",
    "✅ WhatsApp connection (Twilio API)",
    "✅ Debugging & error handling (logging + fallbacks)",
    "✅ Performance & safety checks (validation + rate limits)",
    "✅ Technical documentation (TECHNICAL_DOCUMENTATION.docx)",
    "✅ Final pitch presentation (this PPTX)",
]
for item in core_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Calibri'
    p.space_before = Pt(5)

# Bonus features
bonus_shape = add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, PURPLE)
tf2 = bonus_shape.text_frame
tf2.word_wrap = True
tf2.paragraphs[0].text = "Bonus Features"
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.color.rgb = WHITE
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.name = 'Calibri'

bonus_items = [
    "✅ Computer Vision (Gemini Vision API)",
    "✅ Crisp logic estimator (rule-based, 70+ models)",
    "✅ Modern UI (dark cyberpunk theme)",
    "✅ Session management (multi-turn conversations)",
    "✅ Django Admin (all 5 models registered)",
    "✅ .env configuration (load_dotenv override)",
    "✅ .gitignore + .env.example + README.md",
    "✅ Mock fallback mode (works without API key)",
    "✅ Health check endpoint (/health/)",
]
for item in bonus_items:
    p = tf2.add_paragraph()
    p.text = item
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Calibri'
    p.space_before = Pt(5)

# ================================================================
# SLIDE 13: LIVE DEMO
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0), Inches(0.5), Inches(13.333), Inches(0.8),
             "Live Demo", 36, GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.3), Inches(13.333), Inches(0.5),
             "Open  localhost:8000", 18, CYAN, alignment=PP_ALIGN.CENTER)

demos = [
    ("💬", "1. Chat", '"What\'s the best car for\na family of 5 under 20K JOD?"', PURPLE),
    ("💰", "2. Price Estimator", "Toyota → Camry → 2020\n→ 80,000 km", ORANGE),
    ("📷", "3. Vision", "Upload any car photo →\ninstant AI identification", CYAN),
    ("📲", "4. WhatsApp", 'Click "Send via WhatsApp"\n→ check your phone', GREEN),
]
for i, (icon, title, body, clr) in enumerate(demos):
    col = i % 2
    row = i // 2
    card_with_text(slide, Inches(1.5 + col * 5.5), Inches(2.2 + row * 2.4),
                   Inches(5), Inches(2.0), icon, title, body, clr)

add_text_box(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.4),
             "All features are live and connected to real APIs", 12, DARK_GRAY,
             alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 14: THANK YOU
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0), Inches(0.8), Inches(13.333), Inches(0.8),
             "🚗✨", 60, WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.8), Inches(13.333), Inches(1),
             "Thank You!", 48, CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.9), Inches(13.333), Inches(0.5),
             "IntelliWheels — AI Car Marketplace for Jordan", 18, GRAY, alignment=PP_ALIGN.CENTER)

# Tags
tag_y2 = Inches(3.7)
for i, (txt, clr) in enumerate(tags):
    add_tag(slide, tag_start + Inches(i * 1.7), tag_y2, txt, clr)

# Stats
stats_final = [("7", "API Endpoints", CYAN), ("70+", "Car Models", ORANGE),
               ("4", "AI Features", GREEN), ("5", "DB Models", RED)]
for i, (num, label, clr) in enumerate(stats_final):
    x = Inches(2.5 + i * 2.2)
    add_text_box(slide, x, Inches(4.6), Inches(2), Inches(0.7), num, 36, clr,
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.2), Inches(2), Inches(0.4), label, 12, GRAY,
                 alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.0), Inches(13.333), Inches(0.4),
             "AI Engineer Course — Capstone Project — February 2026", 12, DARK_GRAY,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(0.5),
             "❓ Questions?", 18, GRAY, alignment=PP_ALIGN.CENTER)

# ================================================================
# SAVE
# ================================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'IntelliWheels_Pitch.pptx')
prs.save(output_path)
print(f"Saved: {output_path}")
